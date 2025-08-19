"""
PPTX презентаций экстрактор
Извлекает заголовки, содержимое слайдов и заметки спикера
"""

import logging
from typing import Dict, Any, List

logger = logging.getLogger(__name__)

# Проверяем наличие python-pptx
try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logger.warning("python-pptx не установлен - PPTX обработка недоступна")


def extract_text_from_pptx(path: str) -> Dict[str, Any]:
    """
    Извлекает текст из PPTX презентации.
    
    Args:
        path: путь к PPTX файлу
    
    Returns:
        Dict с полями: success, text, method, meta, error
    """
    
    if not PPTX_AVAILABLE:
        return {
            "success": False, 
            "error": "python-pptx не установлен - установите пакет python-pptx"
        }
    
    try:
        prs = Presentation(path)
    except Exception as e:
        logger.error(f"Ошибка открытия PPTX {path}: {e}")
        return {"success": False, "error": f"Не удалось открыть PPTX: {e}"}

    slides_meta: List[dict] = []
    combined_text: List[str] = []
    total_slides = len(prs.slides)
    
    logger.info(f"Обрабатываю PPTX: {total_slides} слайдов")

    for idx, slide in enumerate(prs.slides, start=1):
        try:
            # Извлекаем заголовок слайда
            title = ""
            if hasattr(slide.shapes, "title") and slide.shapes.title and slide.shapes.title.text:
                title = slide.shapes.title.text.strip()
            
            # Извлекаем текст из всех текстовых блоков
            content_parts: List[str] = []
            
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    shape_text = _extract_text_from_shape(shape)
                    if shape_text and shape_text != title:  # Избегаем дублирования заголовка
                        content_parts.append(shape_text)
            
            # Извлекаем заметки спикера
            notes = ""
            try:
                if hasattr(slide, "notes_slide") and slide.notes_slide:
                    notes_slide = slide.notes_slide
                    if hasattr(notes_slide, "notes_text_frame") and notes_slide.notes_text_frame:
                        notes = (notes_slide.notes_text_frame.text or "").strip()
            except Exception as notes_err:
                logger.debug(f"Ошибка извлечения заметок со слайда {idx}: {notes_err}")
            
            # Формируем метаданные слайда
            slide_meta = {
                "index": idx,
                "title": title,
                "content": content_parts,
                "notes": notes,
                "has_content": bool(title or content_parts or notes)
            }
            slides_meta.append(slide_meta)
            
            # Формируем текстовый блок для слайда
            slide_text_parts = []
            
            if title:
                slide_text_parts.append(f"=== Слайд {idx}: {title} ===")
            else:
                slide_text_parts.append(f"=== Слайд {idx} ===")
            
            if content_parts:
                for part in content_parts:
                    if part.strip():
                        slide_text_parts.append(part.strip())
            
            if notes:
                slide_text_parts.append(f"📝 Заметки спикера: {notes}")
            
            if slide_text_parts:
                combined_text.append("\n".join(slide_text_parts))
                
            logger.debug(f"Слайд {idx}: заголовок='{title[:50]}...', контент={len(content_parts)} блоков, заметки={bool(notes)}")
            
        except Exception as slide_err:
            logger.error(f"Ошибка обработки слайда {idx}: {slide_err}")
            slides_meta.append({
                "index": idx,
                "title": f"Ошибка обработки слайда {idx}",
                "content": [],
                "notes": "",
                "has_content": False,
                "error": str(slide_err)
            })

    # Объединяем весь текст
    full_text = "\n\n".join(combined_text).strip()
    
    if not full_text:
        return {
            "success": False, 
            "error": "Не удалось извлечь текст из PPTX - презентация пустая или не содержит текста"
        }

    # Подсчитываем статистику
    slides_with_content = sum(1 for slide in slides_meta if slide["has_content"])
    slides_with_notes = sum(1 for slide in slides_meta if slide["notes"])
    
    logger.info(f"PPTX обработан: {len(full_text)} символов из {slides_with_content}/{total_slides} слайдов")
    
    return {
        "success": True,
        "text": full_text,
        "method": "python-pptx",
        "meta": {
            "slides": slides_meta,
            "total_slides": total_slides,
            "slides_with_content": slides_with_content,
            "slides_with_notes": slides_with_notes,
            "extraction_stats": f"{slides_with_content} слайдов с контентом, {slides_with_notes} с заметками"
        }
    }


def _extract_text_from_shape(shape) -> str:
    """Извлекает текст из текстового блока слайда"""
    try:
        if not hasattr(shape, "text_frame") or not shape.text_frame:
            return ""
        
        text_parts = []
        
        # Проходим по всем параграфам
        for paragraph in shape.text_frame.paragraphs:
            paragraph_text = paragraph.text.strip()
            if paragraph_text:
                # Определяем уровень списка для форматирования
                level = getattr(paragraph, "level", 0)
                indent = "  " * level
                
                # Добавляем маркеры для списков
                if level > 0:
                    paragraph_text = f"{indent}• {paragraph_text}"
                
                text_parts.append(paragraph_text)
        
        return "\n".join(text_parts).strip()
        
    except Exception as e:
        logger.debug(f"Ошибка извлечения текста из блока: {e}")
        return ""


def extract_pptx_summary_info(slides_meta: List[dict]) -> str:
    """Создает краткую сводку по презентации для пользователя"""
    if not slides_meta:
        return "Пустая презентация"
    
    total_slides = len(slides_meta)
    slides_with_content = sum(1 for slide in slides_meta if slide.get("has_content", False))
    slides_with_notes = sum(1 for slide in slides_meta if slide.get("notes"))
    
    # Собираем заголовки первых нескольких слайдов
    titles = []
    for slide in slides_meta[:5]:  # Первые 5 слайдов
        title = slide.get("title", "").strip()
        if title:
            titles.append(f"• {title}")
        else:
            titles.append(f"• Слайд {slide.get('index', '?')} (без заголовка)")
    
    summary_parts = [
        f"📊 Презентация: {total_slides} слайдов",
        f"📄 Контент найден на {slides_with_content} слайдах"
    ]
    
    if slides_with_notes > 0:
        summary_parts.append(f"📝 Заметки спикера: {slides_with_notes} слайдов")
    
    if titles:
        summary_parts.append("\n🎯 Основные слайды:")
        summary_parts.extend(titles)
        
        if total_slides > 5:
            summary_parts.append(f"• ... и еще {total_slides - 5} слайдов")
    
    return "\n".join(summary_parts)


def check_pptx_availability() -> bool:
    """Проверяет доступность PPTX обработки"""
    return PPTX_AVAILABLE


def get_pptx_info() -> str:
    """Возвращает информацию о поддержке PPTX"""
    if PPTX_AVAILABLE:
        return "PPTX презентации полностью поддерживаются"
    else:
        return "PPTX обработка недоступна - установите python-pptx"