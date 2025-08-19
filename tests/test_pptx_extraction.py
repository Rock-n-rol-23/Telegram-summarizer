#!/usr/bin/env python3
"""
Тесты для PPTX экстрактора
"""

import os
import sys
from pathlib import Path

# Добавляем корневую директорию в путь
sys.path.insert(0, str(Path(__file__).parent.parent))

try:
    from content_extraction.pptx_extractor import extract_text_from_pptx, check_pptx_availability, get_pptx_info, extract_pptx_summary_info
    from pptx import Presentation
    from pptx.util import Inches
except ImportError as e:
    print(f"❌ PPTX зависимости недоступны: {e}")
    print("Установите python-pptx для тестирования PPTX функций")
    sys.exit(1)


def create_test_pptx(output_path: str):
    """Создает тестовую PPTX презентацию"""
    prs = Presentation()
    
    # Слайд 1: Заголовочный слайд
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])  # Title slide layout
    slide1.shapes.title.text = "Тестовая презентация"
    subtitle = slide1.placeholders[1]
    subtitle.text = "Проверка извлечения текста из PPTX"
    
    # Добавляем заметки спикера
    notes_slide1 = slide1.notes_slide
    notes_slide1.notes_text_frame.text = "Это заметки спикера для первого слайда. Здесь может быть дополнительная информация."
    
    # Слайд 2: Слайд с содержимым
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # Content slide layout
    slide2.shapes.title.text = "Основные моменты"
    
    content = slide2.placeholders[1]
    tf = content.text_frame
    tf.text = "Первый пункт презентации"
    
    # Добавляем дополнительные пункты
    p = tf.add_paragraph()
    p.text = "Второй пункт с важной информацией"
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "Подпункт второго пункта"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Третий основной пункт"
    p.level = 0
    
    # Заметки для второго слайда
    notes_slide2 = slide2.notes_slide
    notes_slide2.notes_text_frame.text = "Заметки для слайда с основными моментами. Здесь спикер может найти дополнительные детали."
    
    # Слайд 3: Слайд только с заголовком
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank slide
    
    # Добавляем текстовый блок вручную
    textbox = slide3.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
    text_frame = textbox.text_frame
    text_frame.text = "Заключение"
    
    p = text_frame.add_paragraph()
    p.text = "Это заключительный слайд презентации"
    
    p = text_frame.add_paragraph()
    p.text = "English text: Summary and conclusions"
    
    # Слайд 4: Пустой слайд (для тестирования)
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank slide без контента
    
    # Сохраняем презентацию
    prs.save(output_path)


def test_pptx_availability():
    """Тест проверки доступности PPTX обработки"""
    print("🧪 Тест доступности PPTX")
    
    available = check_pptx_availability()
    print(f"PPTX поддержка: {'✅' if available else '❌'}")
    
    info = get_pptx_info()
    print(f"Статус: {info}")
    
    assert available, "PPTX поддержка должна быть доступна"
    print("✅ Тест доступности прошел")


def test_pptx_extraction():
    """Тест извлечения текста из PPTX"""
    print("\n🧪 Тест извлечения из PPTX")
    
    test_pptx_path = "test_presentation.pptx"
    try:
        # Создаем тестовую презентацию
        create_test_pptx(test_pptx_path)
        assert os.path.exists(test_pptx_path), "Тестовая PPTX не создана"
        
        # Извлекаем текст
        result = extract_text_from_pptx(test_pptx_path)
        
        assert result['success'], f"Извлечение не удалось: {result.get('error')}"
        assert len(result['text']) >= 100, f"Слишком мало текста: {len(result['text'])} символов"
        
        # Проверяем основные элементы текста
        text = result['text']
        assert "Тестовая презентация" in text, "Заголовок презентации не найден"
        assert "Основные моменты" in text, "Заголовок слайда не найден"
        assert "Первый пункт" in text, "Содержимое слайда не найдено"
        assert "заметки спикера" in text.lower(), "Заметки спикера не найдены"
        
        # Проверяем метод и метаданные
        assert result['method'] == 'python-pptx', f"Неверный метод: {result['method']}"
        
        meta = result['meta']
        assert 'slides' in meta, "Метаданные слайдов не найдены"
        assert meta['total_slides'] >= 3, f"Ожидалось минимум 3 слайда, получено {meta['total_slides']}"
        assert meta['slides_with_content'] >= 2, f"Минимум 2 слайда должны иметь контент"
        
        # Проверяем структуру слайдов
        slides = meta['slides']
        assert len(slides) >= 3, f"Ожидалось минимум 3 слайда в метаданных"
        
        # Первый слайд должен иметь заголовок
        slide1 = slides[0]
        assert slide1['title'] == "Тестовая презентация", f"Неверный заголовок первого слайда: {slide1['title']}"
        assert slide1['notes'], "Заметки первого слайда не найдены"
        
        print(f"✅ Извлечено {len(text)} символов из {meta['total_slides']} слайдов")
        print(f"✅ Слайды с контентом: {meta['slides_with_content']}")
        print(f"✅ Слайды с заметками: {meta['slides_with_notes']}")
        
    finally:
        # Очищаем тестовый файл
        if os.path.exists(test_pptx_path):
            os.remove(test_pptx_path)


def test_pptx_summary_info():
    """Тест создания сводки по презентации"""
    print("\n🧪 Тест сводки по презентации")
    
    # Создаем тестовые метаданные слайдов
    test_slides = [
        {
            "index": 1,
            "title": "Введение",
            "content": ["Добро пожаловать", "Цели презентации"],
            "notes": "Заметки для введения",
            "has_content": True
        },
        {
            "index": 2, 
            "title": "Основная часть",
            "content": ["Пункт 1", "Пункт 2", "Пункт 3"],
            "notes": "",
            "has_content": True
        },
        {
            "index": 3,
            "title": "",
            "content": [],
            "notes": "",
            "has_content": False
        }
    ]
    
    summary = extract_pptx_summary_info(test_slides)
    
    assert "3 слайдов" in summary, "Информация о количестве слайдов не найдена"
    assert "2 слайдах" in summary, "Информация о слайдах с контентом не найдена"
    assert "Введение" in summary, "Заголовок первого слайда не найден"
    assert "Основная часть" in summary, "Заголовок второго слайда не найден"
    
    print(f"✅ Сводка создана успешно")
    print(f"Содержимое сводки:\n{summary}")


def test_empty_presentation():
    """Тест обработки пустой презентации"""
    print("\n🧪 Тест пустой презентации")
    
    test_pptx_path = "test_empty.pptx"
    try:
        # Создаем пустую презентацию
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank slide
        prs.save(test_pptx_path)
        
        # Пытаемся извлечь текст
        result = extract_text_from_pptx(test_pptx_path)
        
        assert not result['success'], "Пустая презентация должна возвращать ошибку"
        assert "не содержит текста" in result['error'], f"Неожиданная ошибка: {result['error']}"
        
        print("✅ Пустая презентация корректно обработана")
        
    finally:
        if os.path.exists(test_pptx_path):
            os.remove(test_pptx_path)


def test_error_handling():
    """Тест обработки ошибок"""
    print("\n🧪 Тест обработки ошибок")
    
    # Тест несуществующего файла
    result = extract_text_from_pptx("nonexistent_file.pptx")
    assert not result['success'], "Должна быть ошибка для несуществующего файла"
    assert "Не удалось открыть PPTX" in result['error'], f"Неожиданная ошибка: {result['error']}"
    
    print("✅ Обработка ошибок работает корректно")


def run_all_tests():
    """Запускает все тесты"""
    print("🧪 Тестирование PPTX экстрактора\n")
    
    tests = [
        test_pptx_availability,
        test_pptx_extraction,
        test_pptx_summary_info,
        test_empty_presentation,
        test_error_handling
    ]
    
    passed = 0
    failed = 0
    
    for test in tests:
        try:
            test()
            passed += 1
        except Exception as e:
            print(f"❌ {test.__name__}: {e}")
            failed += 1
    
    print(f"\n📊 Результаты тестов PPTX: {passed} прошли, {failed} провалились")
    return failed == 0


if __name__ == "__main__":
    success = run_all_tests()
    sys.exit(0 if success else 1)